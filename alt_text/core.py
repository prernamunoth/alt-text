from pathlib import Path
from typing import Union, Callable, Dict
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.placeholder import PlaceholderPicture
import click
import shutil
from alt_text.model import AltTextModel

def check_alt_text(pptx_path: Union[str, Path], model: AltTextModel, progress_callback: Callable[[int, int], None] = None) -> Dict:
    """
    Check alt text for all images in a PowerPoint presentation.
    Generate alt text for images without it and save the presentation.
    
    Args:
        pptx_path (Union[str, Path]): Path to the PPTX file
        model (AltTextModel): Model instance to use for alt text generation
        progress_callback (Callable[[int, int], None], optional): Callback function for progress updates
            First argument is current slide number, second is total slides
    
    Returns:
        Dict: Statistics about the processing
    """
    try:
        # Convert string path to Path object if needed
        path = Path(pptx_path) if isinstance(pptx_path, str) else pptx_path
        
        # Create images directory if it doesn't exist
        images_dir = Path("images")
        images_dir.mkdir(exist_ok=True)
        
        # Load the presentation
        prs = Presentation(path)
        
        # Initialize counters
        total_images: int = 0
        images_without_alt: int = 0
        modified = False
        total_slides = len(prs.slides)
        
        # Iterate through all slides
        for slide_number, slide in enumerate(prs.slides, 1):
            if progress_callback:
                progress_callback(slide_number, total_slides)
            
            click.echo(f"\nChecking Slide {slide_number}:")
            
            # Check all shapes in the slide
            for shape in slide.shapes:
                if (shape.shape_type in [
                    MSO_SHAPE_TYPE.PICTURE,
                    MSO_SHAPE_TYPE.MEDIA,
                ] or 
                    isinstance(shape, PlaceholderPicture)):
                    total_images += 1
                    
                    # Save the image and get its path
                    image_path = None
                    try:
                        if hasattr(shape, 'image'):
                            # Use the actual image format extension from python-pptx
                            ext = getattr(shape.image, 'ext', 'png') or 'png'
                            image_filename = f"slide{slide_number}_{shape.name}.{ext}"
                            image_path = images_dir / image_filename
                            
                            with open(image_path, 'wb') as f:
                                f.write(shape.image.blob)
                            click.echo(f"Saved image to: {image_path}")
                    except Exception as e:
                        click.echo(f"Could not save image: {e}")
                        continue
                    
                    # Try different methods to get alt text
                    alt_text: str = ""
                    
                    if hasattr(shape, "alt_text"):
                        alt_text = shape.alt_text or ""
                    
                    if not alt_text and hasattr(shape, "_element"):
                        pic = shape._element
                        if hasattr(pic, "nvPicPr"):
                            props = pic.nvPicPr
                            if hasattr(props, "cNvPr"):
                                alt_text = props.cNvPr.get('descr', "")
                    
                    alt_text = alt_text.strip()
                    
                    click.echo(f"\nImage found in slide {slide_number}:")
                    click.echo(f"Shape type: {shape.shape_type}")
                    click.echo(f"Shape name: {shape.name}")
                    
                    if not alt_text and image_path:
                        images_without_alt += 1
                        click.echo("Generating alt text for this image...")
                        
                        try:
                            # Generate alt text using the model
                            alt_text = model.generate_alt_text(str(image_path))
                            click.echo(f"Generated alt text: '{alt_text}'")
                            
                            # Update the shape's alt text
                            try:
                                # Method 1: Direct alt_text property
                                shape.alt_text = alt_text
                                
                                # Method 2: Through the shape element's properties
                                if hasattr(shape, '_element'):
                                    pic = shape._element
                                    if hasattr(pic, 'nvPicPr'):
                                        props = pic.nvPicPr
                                        if hasattr(props, 'cNvPr'):
                                            props.cNvPr.set('descr', alt_text)
                                
                                modified = True
                                click.echo("Successfully set alt text")
                                
                            except Exception as e:
                                click.echo(f"Warning: Issue setting alt text: {e}")
                            
                        except Exception as e:
                            click.echo(f"Error generating alt text: {e}")
                    else:
                        click.echo(f"Existing alt text: '{alt_text}'")
        
        # Save the presentation if modifications were made
        if modified:
            output_path = path.parent / f"updated_{path.name}"
            try:
                prs.save(output_path)
                
                if output_path.exists():
                    click.echo(f"\nSaved updated presentation to: {output_path}")
                else:
                    click.echo("Warning: File save may have failed")
            except Exception as save_error:
                click.echo(f"Error saving presentation: {save_error}")
        
        # Return statistics
        return {
            "total_slides": total_slides,
            "total_images": total_images,
            "images_without_alt": images_without_alt,
            "images_with_alt": total_images - images_without_alt,
            "modified": modified,
            "output_path": str(output_path) if modified else None
        }
        
    except Exception as e:
        click.echo(f"Error processing file: {e}")
        return None

@click.command()
@click.argument('pptx_path', type=click.Path(exists=True, dir_okay=False))
@click.option('--output', '-o', type=click.Path(dir_okay=False), help='Output file path (optional)')
def main(pptx_path: str, output: str = None) -> None:
    """
    Process a PowerPoint file to add alt text to images.
    
    PPTX_PATH: Path to the PowerPoint file to process
    """
    try:
        # Process the file
        stats = check_alt_text(pptx_path)
        
        if stats and stats["modified"]:
            # If output path is specified, copy the file there
            if output:
                input_path = Path(pptx_path)
                output_path = Path(output)
                if output_path.parent != input_path.parent:
                    output_path.parent.mkdir(parents=True, exist_ok=True)
                shutil.copy2(input_path.parent / f"updated_{input_path.name}", output_path)
                click.echo(f"\nCopied updated file to: {output_path}")
            
            # Print summary
            click.echo("\n=== Summary ===")
            click.echo(f"Total slides: {stats['total_slides']}")
            click.echo(f"Total images found: {stats['total_images']}")
            click.echo(f"Images without alt text: {stats['images_without_alt']}")
            click.echo(f"Images with alt text: {stats['images_with_alt']}")
            
    except Exception as e:
        click.echo(f"Error: {e}")
        raise click.Abort()

if __name__ == "__main__":
    main()